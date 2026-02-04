"""Tests for beamer2animate.pptx_builder."""

import os

import pytest
from pptx.util import Inches, Pt

from beamer2animate.pptx_builder import (
    create_presentation,
    add_slide,
    add_title_to_slide,
    add_image_to_slide,
    add_appear_animation,
    add_disappear_animation,
    save_presentation,
    NSMAP,
)


class TestCreatePresentation:
    def test_default_dimensions(self):
        prs = create_presentation()
        assert prs.slide_width == Inches(13.333)
        assert prs.slide_height == Inches(7.5)

    def test_custom_dimensions(self):
        prs = create_presentation(width=10, height=7.5)
        assert prs.slide_width == Inches(10)
        assert prs.slide_height == Inches(7.5)


class TestAddSlide:
    def test_adds_slide(self):
        prs = create_presentation()
        slide = add_slide(prs)
        assert len(prs.slides) == 1

    def test_multiple_slides(self):
        prs = create_presentation()
        add_slide(prs)
        add_slide(prs)
        add_slide(prs)
        assert len(prs.slides) == 3


class TestAddTitle:
    def test_title_text(self):
        prs = create_presentation()
        slide = add_slide(prs)
        tb = add_title_to_slide(slide, "My Title")
        assert tb.text_frame.paragraphs[0].text == "My Title"

    def test_title_font(self):
        prs = create_presentation()
        slide = add_slide(prs)
        tb = add_title_to_slide(slide, "Title")
        p = tb.text_frame.paragraphs[0]
        assert p.font.size == Pt(32)
        assert p.font.bold is True

    def test_title_position(self):
        prs = create_presentation()
        slide = add_slide(prs)
        tb = add_title_to_slide(slide, "Title")
        assert tb.left == Inches(1.0)
        assert tb.top == Inches(0.3)
        assert tb.width == Inches(11.333)
        assert tb.height == Inches(0.8)

    def test_custom_font_size(self):
        prs = create_presentation()
        slide = add_slide(prs)
        tb = add_title_to_slide(slide, "Big", font_size=48)
        assert tb.text_frame.paragraphs[0].font.size == Pt(48)


class TestAddImage:
    @pytest.fixture
    def small_png(self, tmp_dir):
        """Create a small test PNG."""
        from PIL import Image
        path = os.path.join(tmp_dir, "test.png")
        img = Image.new("RGB", (400, 200), color="red")
        img.save(path)
        return path

    def test_adds_image(self, small_png):
        prs = create_presentation()
        slide = add_slide(prs)
        shape = add_image_to_slide(slide, small_png, left=1.0, top=1.0, width=5.0)
        assert shape is not None
        assert shape.width == Inches(5.0)

    def test_aspect_ratio_preserved(self, small_png):
        prs = create_presentation()
        slide = add_slide(prs)
        # Image is 400x200 = 2:1 aspect ratio
        shape = add_image_to_slide(slide, small_png, left=0, top=0, width=10.0)
        # Height should be 5.0 inches (half of width due to 2:1 ratio)
        assert abs(shape.height - Inches(5.0)) < Inches(0.01)

    def test_position(self, small_png):
        prs = create_presentation()
        slide = add_slide(prs)
        shape = add_image_to_slide(slide, small_png, left=2.0, top=3.0, width=4.0)
        assert shape.left == Inches(2.0)
        assert shape.top == Inches(3.0)


class TestAnimations:
    @pytest.fixture
    def slide_with_shapes(self, tmp_dir):
        """Create a slide with two image shapes for animation testing."""
        from PIL import Image

        prs = create_presentation()
        slide = add_slide(prs)
        shapes = []
        for i in range(3):
            path = os.path.join(tmp_dir, f"shape_{i}.png")
            img = Image.new("RGB", (200, 100), color="blue")
            img.save(path)
            shape = add_image_to_slide(slide, path, left=0, top=i * 2.0, width=5.0)
            shapes.append(shape)
        return slide, shapes

    def test_appear_animation_creates_timing(self, slide_with_shapes):
        slide, shapes = slide_with_shapes
        add_appear_animation(slide, shapes[0], click_index=0)

        timing = slide._element.find('.//p:timing', NSMAP)
        assert timing is not None

    def test_appear_animation_targets_shape(self, slide_with_shapes):
        slide, shapes = slide_with_shapes
        add_appear_animation(slide, shapes[0], click_index=0)

        spTgt = slide._element.find('.//p:spTgt', NSMAP)
        assert spTgt is not None
        assert spTgt.get('spid') == str(shapes[0].shape_id)

    def test_multiple_appear_animations(self, slide_with_shapes):
        slide, shapes = slide_with_shapes
        for i, shape in enumerate(shapes):
            add_appear_animation(slide, shape, click_index=i)

        # Should have animation targets for all shapes
        spTgts = slide._element.findall('.//p:spTgt', NSMAP)
        targeted_ids = {t.get('spid') for t in spTgts}
        for shape in shapes:
            assert str(shape.shape_id) in targeted_ids

    def test_disappear_animation(self, slide_with_shapes):
        slide, shapes = slide_with_shapes
        # Need appear first to set up timing structure
        add_appear_animation(slide, shapes[0], click_index=0)
        add_appear_animation(slide, shapes[1], click_index=1)
        add_disappear_animation(slide, shapes[0], click_index=1)

        # Check that an exit animation exists
        exit_nodes = slide._element.findall('.//' + '{%s}cTn[@presetClass="exit"]' % NSMAP['p'])
        assert len(exit_nodes) > 0


class TestSavePresentation:
    def test_save(self, tmp_dir):
        prs = create_presentation()
        add_slide(prs)
        path = os.path.join(tmp_dir, "test.pptx")
        save_presentation(prs, path)
        assert os.path.exists(path)
        assert os.path.getsize(path) > 0
