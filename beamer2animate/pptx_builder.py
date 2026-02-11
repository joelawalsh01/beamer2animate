"""Create PowerPoint presentations with animations."""

import os
from typing import List, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from lxml import etree

# PowerPoint XML namespaces
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}


def create_presentation(width: float = 13.333, height: float = 7.5) -> Presentation:
    """Create a new presentation with 16:9 aspect ratio."""
    prs = Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)
    return prs


def add_slide(prs: Presentation, layout_index: int = 6) -> 'Slide':
    """Add a blank slide to the presentation.

    layout_index 6 is typically the blank layout.
    """
    slide_layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(slide_layout)


def add_title_to_slide(slide, title: str, font_size: int = 32):
    """Add a title text box to the slide."""
    left = Inches(1.0)
    top = Inches(0.3)
    width = Inches(11.333)
    height = Inches(0.8)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    return textbox


def add_image_to_slide(
    slide,
    image_path: str,
    left: float,
    top: float,
    width: Optional[float] = None,
    height: Optional[float] = None
) -> 'Shape':
    """Add an image to the slide.

    Args:
        slide: The slide to add the image to
        image_path: Path to the image file
        left: Left position in inches
        top: Top position in inches
        width: Width in inches (optional, maintains aspect ratio if only one specified)
        height: Height in inches (optional)

    Returns:
        The picture shape
    """
    left_emu = Inches(left)
    top_emu = Inches(top)
    width_emu = Inches(width) if width else None
    height_emu = Inches(height) if height else None

    if width_emu and not height_emu:
        picture = slide.shapes.add_picture(image_path, left_emu, top_emu, width=width_emu)
    elif height_emu and not width_emu:
        picture = slide.shapes.add_picture(image_path, left_emu, top_emu, height=height_emu)
    elif width_emu and height_emu:
        picture = slide.shapes.add_picture(image_path, left_emu, top_emu, width=width_emu, height=height_emu)
    else:
        picture = slide.shapes.add_picture(image_path, left_emu, top_emu)

    return picture


def add_content_background(slide, top=0.74, height=6.11):
    """Add a white rectangle covering the Beamer content area.

    This hides the static content from the Beamer PDF background while
    preserving the theme chrome (header banner with frame title, footer).
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(top),
        Inches(13.333), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return shape


def get_shape_id(shape) -> int:
    """Get the shape ID for animation targeting."""
    return shape.shape_id


def add_appear_animation(slide, shape, click_index: int = 0, delay: int = 0):
    """Add an 'appear' animation to a shape.

    This makes the shape invisible initially and appear on click.

    Args:
        slide: The slide containing the shape
        shape: The shape to animate
        click_index: Which click triggers this animation (0 = first click)
        delay: Delay in milliseconds before animation starts
    """
    shape_id = get_shape_id(shape)

    # Get or create the timing element
    timing = slide._element.find('.//p:timing', NSMAP)
    if timing is None:
        timing = etree.SubElement(slide._element, '{%s}timing' % NSMAP['p'])

    # Get or create tnLst (time node list)
    tnLst = timing.find('p:tnLst', NSMAP)
    if tnLst is None:
        tnLst = etree.SubElement(timing, '{%s}tnLst' % NSMAP['p'])

    # Get or create the main parallel time node
    main_par = tnLst.find('p:par', NSMAP)
    if main_par is None:
        main_par = etree.SubElement(tnLst, '{%s}par' % NSMAP['p'])
        main_cTn = etree.SubElement(main_par, '{%s}cTn' % NSMAP['p'])
        main_cTn.set('id', '1')
        main_cTn.set('dur', 'indefinite')
        main_cTn.set('restart', 'never')
        main_cTn.set('nodeType', 'tmRoot')

    main_cTn = main_par.find('p:cTn', NSMAP)
    if main_cTn is None:
        main_cTn = etree.SubElement(main_par, '{%s}cTn' % NSMAP['p'])
        main_cTn.set('id', '1')
        main_cTn.set('dur', 'indefinite')
        main_cTn.set('restart', 'never')
        main_cTn.set('nodeType', 'tmRoot')

    # Get or create childTnLst
    childTnLst = main_cTn.find('p:childTnLst', NSMAP)
    if childTnLst is None:
        childTnLst = etree.SubElement(main_cTn, '{%s}childTnLst' % NSMAP['p'])

    # Find or create the sequence for this click
    seq = None
    for existing_seq in childTnLst.findall('p:seq', NSMAP):
        # Check if this sequence is for our click index
        seq = existing_seq
        break

    if seq is None:
        seq = etree.SubElement(childTnLst, '{%s}seq' % NSMAP['p'])
        seq.set('concurrent', '1')
        seq.set('nextAc', 'seek')

        seq_cTn = etree.SubElement(seq, '{%s}cTn' % NSMAP['p'])
        seq_cTn.set('id', '2')
        seq_cTn.set('dur', 'indefinite')
        seq_cTn.set('nodeType', 'mainSeq')

        seq_childTnLst = etree.SubElement(seq_cTn, '{%s}childTnLst' % NSMAP['p'])

        # Add prevCondLst for click trigger
        prevCondLst = etree.SubElement(seq, '{%s}prevCondLst' % NSMAP['p'])
        cond = etree.SubElement(prevCondLst, '{%s}cond' % NSMAP['p'])
        cond.set('evt', 'onPrev')
        cond.set('delay', '0')
        tgtEl = etree.SubElement(cond, '{%s}tgtEl' % NSMAP['p'])
        sldTgt = etree.SubElement(tgtEl, '{%s}sldTgt' % NSMAP['p'])

        # Add nextCondLst
        nextCondLst = etree.SubElement(seq, '{%s}nextCondLst' % NSMAP['p'])
        cond = etree.SubElement(nextCondLst, '{%s}cond' % NSMAP['p'])
        cond.set('evt', 'onNext')
        cond.set('delay', '0')
        tgtEl = etree.SubElement(cond, '{%s}tgtEl' % NSMAP['p'])
        sldTgt = etree.SubElement(tgtEl, '{%s}sldTgt' % NSMAP['p'])
    else:
        seq_cTn = seq.find('p:cTn', NSMAP)
        seq_childTnLst = seq_cTn.find('p:childTnLst', NSMAP)

    # Create the animation for this shape
    # Each click creates a new parallel container
    click_par = etree.SubElement(seq_childTnLst, '{%s}par' % NSMAP['p'])
    click_cTn = etree.SubElement(click_par, '{%s}cTn' % NSMAP['p'])
    click_cTn.set('id', str(100 + click_index * 10))
    click_cTn.set('fill', 'hold')

    # Start condition
    stCondLst = etree.SubElement(click_cTn, '{%s}stCondLst' % NSMAP['p'])
    cond = etree.SubElement(stCondLst, '{%s}cond' % NSMAP['p'])
    cond.set('delay', 'indefinite')

    click_childTnLst = etree.SubElement(click_cTn, '{%s}childTnLst' % NSMAP['p'])

    # Inner parallel for the actual animation
    anim_par = etree.SubElement(click_childTnLst, '{%s}par' % NSMAP['p'])
    anim_cTn = etree.SubElement(anim_par, '{%s}cTn' % NSMAP['p'])
    anim_cTn.set('id', str(101 + click_index * 10))
    anim_cTn.set('fill', 'hold')

    anim_stCondLst = etree.SubElement(anim_cTn, '{%s}stCondLst' % NSMAP['p'])
    anim_cond = etree.SubElement(anim_stCondLst, '{%s}cond' % NSMAP['p'])
    anim_cond.set('delay', str(delay))

    anim_childTnLst = etree.SubElement(anim_cTn, '{%s}childTnLst' % NSMAP['p'])

    # Another parallel for the set animation
    set_par = etree.SubElement(anim_childTnLst, '{%s}par' % NSMAP['p'])
    set_cTn = etree.SubElement(set_par, '{%s}cTn' % NSMAP['p'])
    set_cTn.set('id', str(102 + click_index * 10))
    set_cTn.set('presetID', '1')  # 1 = Appear
    set_cTn.set('presetClass', 'entr')  # entrance effect
    set_cTn.set('presetSubtype', '0')
    set_cTn.set('fill', 'hold')
    set_cTn.set('nodeType', 'clickEffect')

    set_stCondLst = etree.SubElement(set_cTn, '{%s}stCondLst' % NSMAP['p'])
    set_cond = etree.SubElement(set_stCondLst, '{%s}cond' % NSMAP['p'])
    set_cond.set('delay', '0')

    set_childTnLst = etree.SubElement(set_cTn, '{%s}childTnLst' % NSMAP['p'])

    # The actual 'set' animation (instant visibility change)
    set_elem = etree.SubElement(set_childTnLst, '{%s}set' % NSMAP['p'])
    set_cBhvr = etree.SubElement(set_elem, '{%s}cBhvr' % NSMAP['p'])
    set_inner_cTn = etree.SubElement(set_cBhvr, '{%s}cTn' % NSMAP['p'])
    set_inner_cTn.set('id', str(103 + click_index * 10))
    set_inner_cTn.set('dur', '1')
    set_inner_cTn.set('fill', 'hold')

    set_inner_stCondLst = etree.SubElement(set_inner_cTn, '{%s}stCondLst' % NSMAP['p'])
    set_inner_cond = etree.SubElement(set_inner_stCondLst, '{%s}cond' % NSMAP['p'])
    set_inner_cond.set('delay', '0')

    # Target element (the shape)
    tgtEl = etree.SubElement(set_cBhvr, '{%s}tgtEl' % NSMAP['p'])
    spTgt = etree.SubElement(tgtEl, '{%s}spTgt' % NSMAP['p'])
    spTgt.set('spid', str(shape_id))

    # Attribute to animate (visibility)
    attrNameLst = etree.SubElement(set_cBhvr, '{%s}attrNameLst' % NSMAP['p'])
    attrName = etree.SubElement(attrNameLst, '{%s}attrName' % NSMAP['p'])
    attrName.text = 'style.visibility'

    # To value (visible)
    to = etree.SubElement(set_elem, '{%s}to' % NSMAP['p'])
    strVal = etree.SubElement(to, '{%s}strVal' % NSMAP['p'])
    strVal.set('val', 'visible')


def add_disappear_animation(slide, shape, click_index: int):
    """Add a 'disappear' animation to a shape (for hiding previous cumulative image)."""
    shape_id = get_shape_id(shape)

    timing = slide._element.find('.//p:timing', NSMAP)
    if timing is None:
        return

    tnLst = timing.find('p:tnLst', NSMAP)
    if tnLst is None:
        return

    main_par = tnLst.find('p:par', NSMAP)
    if main_par is None:
        return

    main_cTn = main_par.find('p:cTn', NSMAP)
    childTnLst = main_cTn.find('p:childTnLst', NSMAP)
    seq = childTnLst.find('p:seq', NSMAP)
    seq_cTn = seq.find('p:cTn', NSMAP)
    seq_childTnLst = seq_cTn.find('p:childTnLst', NSMAP)

    # Find the click container for this click index
    click_pars = seq_childTnLst.findall('p:par', NSMAP)
    if click_index < len(click_pars):
        click_par = click_pars[click_index]
    else:
        return

    click_cTn = click_par.find('p:cTn', NSMAP)
    click_childTnLst = click_cTn.find('p:childTnLst', NSMAP)

    # Add disappear animation to the same click
    anim_par = etree.SubElement(click_childTnLst, '{%s}par' % NSMAP['p'])
    anim_cTn = etree.SubElement(anim_par, '{%s}cTn' % NSMAP['p'])
    anim_cTn.set('id', str(200 + click_index * 10))
    anim_cTn.set('fill', 'hold')

    anim_stCondLst = etree.SubElement(anim_cTn, '{%s}stCondLst' % NSMAP['p'])
    anim_cond = etree.SubElement(anim_stCondLst, '{%s}cond' % NSMAP['p'])
    anim_cond.set('delay', '0')

    anim_childTnLst = etree.SubElement(anim_cTn, '{%s}childTnLst' % NSMAP['p'])

    set_par = etree.SubElement(anim_childTnLst, '{%s}par' % NSMAP['p'])
    set_cTn = etree.SubElement(set_par, '{%s}cTn' % NSMAP['p'])
    set_cTn.set('id', str(201 + click_index * 10))
    set_cTn.set('presetID', '1')
    set_cTn.set('presetClass', 'exit')  # exit effect
    set_cTn.set('presetSubtype', '0')
    set_cTn.set('fill', 'hold')
    set_cTn.set('nodeType', 'withEffect')  # happens with the appear

    set_stCondLst = etree.SubElement(set_cTn, '{%s}stCondLst' % NSMAP['p'])
    set_cond = etree.SubElement(set_stCondLst, '{%s}cond' % NSMAP['p'])
    set_cond.set('delay', '0')

    set_childTnLst = etree.SubElement(set_cTn, '{%s}childTnLst' % NSMAP['p'])

    set_elem = etree.SubElement(set_childTnLst, '{%s}set' % NSMAP['p'])
    set_cBhvr = etree.SubElement(set_elem, '{%s}cBhvr' % NSMAP['p'])
    set_inner_cTn = etree.SubElement(set_cBhvr, '{%s}cTn' % NSMAP['p'])
    set_inner_cTn.set('id', str(202 + click_index * 10))
    set_inner_cTn.set('dur', '1')
    set_inner_cTn.set('fill', 'hold')

    set_inner_stCondLst = etree.SubElement(set_inner_cTn, '{%s}stCondLst' % NSMAP['p'])
    set_inner_cond = etree.SubElement(set_inner_stCondLst, '{%s}cond' % NSMAP['p'])
    set_inner_cond.set('delay', '0')

    tgtEl = etree.SubElement(set_cBhvr, '{%s}tgtEl' % NSMAP['p'])
    spTgt = etree.SubElement(tgtEl, '{%s}spTgt' % NSMAP['p'])
    spTgt.set('spid', str(shape_id))

    attrNameLst = etree.SubElement(set_cBhvr, '{%s}attrNameLst' % NSMAP['p'])
    attrName = etree.SubElement(attrNameLst, '{%s}attrName' % NSMAP['p'])
    attrName.text = 'style.visibility'

    to = etree.SubElement(set_elem, '{%s}to' % NSMAP['p'])
    strVal = etree.SubElement(to, '{%s}strVal' % NSMAP['p'])
    strVal.set('val', 'hidden')


def set_shape_initially_hidden(shape):
    """Set a shape to be initially hidden (for animation)."""
    # Access the shape's spPr (shape properties)
    sp = shape._element
    nvSpPr = sp.find('.//p:nvSpPr', NSMAP) or sp.find('.//p:nvPicPr', NSMAP)
    if nvSpPr is not None:
        cNvPr = nvSpPr.find('p:cNvPr', NSMAP)
        if cNvPr is not None:
            # We can't directly set visibility, but we'll handle it through animations
            pass


def create_pptx(
    title: str = "Presentation",
    output_path: str = "output.pptx"
) -> Presentation:
    """Create a new PowerPoint presentation.

    Args:
        title: Presentation title
        output_path: Where to save the file

    Returns:
        The Presentation object
    """
    prs = create_presentation()
    return prs


def save_presentation(prs: Presentation, output_path: str):
    """Save the presentation to a file."""
    prs.save(output_path)
